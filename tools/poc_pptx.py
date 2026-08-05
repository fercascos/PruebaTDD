#!/usr/bin/env python3
"""Prueba de concepto del bloque 4 · generación de PPTX.

Responde a las preguntas de `docs/15` §21.3 con **números medidos**, no con
opiniones. Genera un informe a partir de la plantilla real del cliente y deja
el resultado listo para renderizar y comparar.

    python3 tools/poc_pptx.py <plantilla.pptx> <salida.pptx>

La plantilla es material confidencial del cliente y **no está en el
repositorio**: se pasa por ruta.
"""

from __future__ import annotations

import sys
import time
from decimal import Decimal
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent / "apps/api/src"))

from pptx import Presentation  # noqa: E402
from tdd.exports.capex_xlsx import generar_xlsx  # noqa: E402
from tdd.reporting import capex_layout as cl  # noqa: E402
from tdd.reporting.clone import clonar_diapositiva, sustituir_marcadores  # noqa: E402
from tdd.reporting.fonts import comprobar_familias  # noqa: E402
from tdd.reporting.overflow import capacidad_del_marco, evaluar  # noqa: E402
from tdd.reporting.pptx_table import insertar_tabla  # noqa: E402
from tdd.reporting.watermark import retirar_marcas_de_agua  # noqa: E402

SISTEMAS = [
    ("Cimentación", "Se aprecian fisuras de retracción en la solera de la zona de aparcamiento, "
     "con anchura inferior a 1 mm y sin continuidad en los muros de contención. No se observan "
     "asientos diferenciales ni signos de agotamiento estructural.", "Estado aceptable."),
    ("Cubierta", "La impermeabilización de la cubierta invertida presenta acumulación de agua en "
     "torno a tres sumideros, con vegetación incipiente en las juntas perimetrales. La lámina "
     "supera su vida útil estimada.", "Se recomienda renovación integral a medio plazo."),
    ("Fachadas", "El revestimiento de panel composite muestra desprendimientos puntuales en la "
     "fachada norte y sellados degradados en los encuentros con carpintería.", "Reparación a corto plazo."),
]

LINEAS = [
    cl.LineaCapex("", "Aparcamiento", "Vida útil", "Analizar estado del subsuelo", "Alto",
                  "Existen fisuras y grietas en la última losa de hormigón y en los muros de "
                  "contención", "CORTO", Decimal("5500.00")),
    cl.LineaCapex("", "Cubierta", "Mantenimiento", "Limpieza de lucernarios", "Bajo",
                  "Gasto operativo. Durante la visita se aprecia que necesita una limpieza en "
                  "profundidad", "CORTO", Decimal("2300.00")),
    cl.LineaCapex("", "Cubierta", "Vida útil", "Renovación de impermeabilización", "Moderado",
                  "Fin de vida útil. La impermeabilización data de la construcción original",
                  "MEDIO", Decimal("83407.50")),
    cl.LineaCapex("", "Fachadas", "Reparación", "Sellado de encuentros con carpintería", "Alto",
                  "Se aprecian entradas de agua puntuales", "CORTO", Decimal("32970.00")),
    cl.LineaCapex("", "Fachadas", "Mejora", "Sustitución de panel composite", "Bajo",
                  "Mejora estética y de eficiencia energética", "MEJORAS", Decimal("294000.00")),
    cl.LineaCapex("", "General", "Normativa", "Adecuación RIPCI", "Extremo",
                  "Incumplimiento del reglamento de instalaciones de protección contra incendios",
                  "CORTO", Decimal("144780.00")),
    cl.LineaCapex("", "General", "Otro", "Petición específica del cliente", "Bajo",
                  "Estudio de viabilidad solicitado fuera del alcance habitual",
                  "OTRO", Decimal("12000.00")),
]


def main() -> int:
    if len(sys.argv) < 3:
        print(__doc__)
        return 2
    plantilla, salida = Path(sys.argv[1]), Path(sys.argv[2])

    print("=" * 78)
    print("PRUEBA DE CONCEPTO · BLOQUE 4 (PPTX)")
    print("=" * 78)

    # ── 1 · Fuentes ─────────────────────────────────────────────────────────
    estado = comprobar_familias()
    faltan = [f for f, ok in estado.items() if not ok]
    print(f"\n1. FUENTES · {len(estado) - len(faltan)} de {len(estado)} familias instaladas")
    if faltan:
        print(f"   ⚠ faltan: {', '.join(faltan)}")

    # ── 2 · Capacidad medida de los marcos ──────────────────────────────────
    print("\n2. CAPACIDAD DE LOS MARCOS (medida con la fuente real)")
    cap_sistema = capacidad_del_marco(
        ancho_in=8.79, alto_in=5.90, cuerpo_pt=10, familia="Gotham Light"
    )
    print(f"   Diapositiva de sistema: {cap_sistema.caracteres_por_linea} car/línea × "
          f"{cap_sistema.lineas} líneas = {cap_sistema.caracteres} caracteres")
    cap_titular = capacidad_del_marco(
        ancho_in=8.79, alto_in=0.40, cuerpo_pt=24, familia="Gotham Ultra",
        muestra="SISTEMA DE CLIMATIZACIÓN Y VENTILACIÓN",
    )
    print(f"   Titular a 24 pt:        {cap_titular.caracteres_por_linea} caracteres por línea")

    for nombre, desc, val in SISTEMAS:
        a = evaluar(f"{desc}\n{val}", cap_sistema)
        print(f"   · {nombre:14s} {len(desc) + len(val):5d} car. → {a.severidad} ({a.ocupacion:.0%})")

    # ── 3 · Clonado ─────────────────────────────────────────────────────────
    t0 = time.monotonic()
    prs = Presentation(str(plantilla))
    n_original = len(prs.slides)
    print(f"\n3. CLONADO · plantilla de {n_original} diapositivas")

    origen = prs.slides[12]  # diapositiva 13: Cimentación + Estructura
    for nombre, desc, val in SISTEMAS:
        nueva = clonar_diapositiva(prs, origen)
        sin_resolver = sustituir_marcadores(
            nueva, {"system.name": nombre, "system.description": desc, "system.assessment": val}
        )
        print(f"   · clonada para «{nombre}» · marcadores sin resolver: {len(sin_resolver)}")

    # ── 4 · Tabla nativa de CAPEX ───────────────────────────────────────────
    print("\n4. TABLA NATIVA DE CAPEX (P-31) Y MARCA DE AGUA (P-43)")
    layout = cl.construir(LINEAS, capitulo="Arquitectura", locale="es-ES")
    print(f"   Columnas: {len(layout.columnas)} · ancho total {layout.ancho_total_in} in")
    print(f"   Filas:    {len(layout.filas)} (secciones + datos + total)")
    print("   Totales:  " + " · ".join(
        f"{k}={cl.formatear_importe(v)}" for k, v in layout.totales.items() if v
    ))

    for trozo in cl.particionar(layout, filas_por_diapositiva=18):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        insertar_tabla(slide, trozo)
    print(f"   Diapositivas de tabla generadas: "
          f"{len(cl.particionar(layout, filas_por_diapositiva=18))}")

    # [REQ] P-43 · La marca de agua no aparece en las versiones generadas.
    marcas = retirar_marcas_de_agua(prs)
    for m in marcas:
        print(f"   · marca de agua retirada del {m.donde}: «{m.texto}»")

    prs.save(str(salida))
    ms = (time.monotonic() - t0) * 1000
    print(f"\n5. RESULTADO · {len(prs.slides)} diapositivas ({n_original} + "
          f"{len(prs.slides) - n_original}) en {ms:.0f} ms → {salida.name}")

    # ── 6 · XLSX desde el MISMO layout ──────────────────────────────────────
    xlsx = salida.with_suffix(".xlsx")
    xlsx.write_bytes(generar_xlsx(layout))
    print(f"6. XLSX desde el mismo CapexTableLayout → {xlsx.name} "
          f"({xlsx.stat().st_size // 1024} KB)")

    # ── 7 · El original, intacto ────────────────────────────────────────────
    import hashlib

    h = hashlib.sha256(plantilla.read_bytes()).hexdigest()
    print(f"7. ORIGINAL INTACTO · sha256 {h[:16]}…")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
