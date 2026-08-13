"""Bloque 4 punta a punta: plantilla, mapeo, avisos, generación y emisión.

`[REQ]` **No se usa la plantilla real del cliente.** Es material confidencial y
no está en el repositorio: aquí se construye una plantilla mínima con
`python-pptx`, con sus marcadores y su marca de agua, que es todo lo que hace
falta para comprobar el comportamiento.

Las dos pruebas que sostienen el bloque son la del **snapshot** —el informe
sigue diciendo lo mismo aunque los datos cambien después— y la de **inmutable
al emitir**, escrita saltándose la API.
"""

from __future__ import annotations

import io
import uuid
from typing import Any

import pytest
from fastapi.testclient import TestClient
from pptx import Presentation
from pptx.util import Inches, Pt
from sqlalchemy import Engine, text

from tests.unit.test_imagenes import imagen

pytestmark = pytest.mark.db

RUTA = "/api/v1"


def plantilla_pptx(*, con_marca_de_agua: bool = True) -> bytes:
    """Una plantilla mínima: portada con marcadores y marca de agua DRAFT."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    caja = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
    marco = caja.text_frame
    marco.text = "Informe {{project.code}} — {{project.name}}"
    marco.add_paragraph().text = "Cliente: {{project.client}}"
    marco.add_paragraph().text = "CAPEX total: {{capex.total}}"

    if con_marca_de_agua:
        auxiliar = prs.slides.add_slide(prs.slide_layouts[6])
        forma = auxiliar.shapes.add_textbox(Inches(2), Inches(2), Inches(6), Inches(2))
        forma.rotation = 315.0
        run = forma.text_frame.paragraphs[0].add_run()
        run.text = "DRAFT"
        run.font.size = Pt(96)
        elemento = forma._element
        elemento.getparent().remove(elemento)
        prs.slide_masters[0].shapes._spTree.append(elemento)

    salida = io.BytesIO()
    prs.save(salida)
    return salida.getvalue()


MAPEO_COMPLETO = {
    "project.code": "project.code",
    "project.name": "project.name",
    "project.client": "project.client",
    "capex.total": "capex.total",
}


@pytest.fixture
def proyecto(cliente: TestClient, cab: Any, datos_base: dict[str, uuid.UUID]) -> str:
    r = cliente.post(
        f"{RUTA}/projects",
        headers=cab("admin_a"),
        json={
            "client_id": str(datos_base["cliente_a"]),
            "internal_code": f"INF-{uuid.uuid4().hex[:6]}",
            "name": "Encargo con informe",
            "applicable_phases": [{"code": "FULL_REPORT"}],
        },
    )
    return str(r.json()["id"])


@pytest.fixture
def plantilla(cliente: TestClient, cab: Any) -> dict[str, Any]:
    r = cliente.post(
        f"{RUTA}/report-templates",
        headers=cab("admin_a"),
        files={
            "file": ("modelo.pptx", io.BytesIO(plantilla_pptx()), "application/vnd.ms-powerpoint")
        },
        data={"name": f"Modelo {uuid.uuid4().hex[:6]}", "language": "es"},
    )
    assert r.status_code == 201, r.text
    return r.json()


@pytest.fixture
def mapeo(cliente: TestClient, cab: Any, plantilla: dict[str, Any]) -> dict[str, Any]:
    r = cliente.post(
        f"{RUTA}/report-templates/{plantilla['id']}/mappings",
        headers=cab("admin_a"),
        json={"name": "Estándar", "bindings": MAPEO_COMPLETO, "is_default": True},
    )
    assert r.status_code == 201, r.text
    return r.json()


@pytest.fixture
def con_hallazgo(
    cliente: TestClient, cab: Any, proyecto: str, motor_admin: Engine
) -> dict[str, Any]:
    """Un proyecto con un activo, un hallazgo validado y su línea de CAPEX."""
    with motor_admin.begin() as conn:
        tipologia = conn.execute(
            text("SELECT id FROM asset_typology ORDER BY sort_order LIMIT 1")
        ).scalar_one()
        zona = conn.execute(
            text(
                "SELECT z.id FROM zone z JOIN zone_typology zt ON zt.zone_id = z.id "
                "WHERE zt.typology_id = :t LIMIT 1"
            ),
            {"t": tipologia},
        ).scalar_one()
        codigo = conn.execute(
            text("SELECT id FROM capex_code WHERE level = 3 LIMIT 1")
        ).scalar_one()

    activo = cliente.post(
        f"{RUTA}/projects/{proyecto}/assets",
        headers=cab("consultor_a"),
        json={"name": "Nave A", "typology_id": str(tipologia)},
    ).json()

    hallazgo = cliente.post(
        f"{RUTA}/projects/{proyecto}/findings",
        headers=cab("consultor_a"),
        json={
            "asset_id": activo["id"],
            "capex_code_id": str(codigo),
            "zone_id": str(zona),
            "title": "Renovación de impermeabilización",
            "description": "La lámina supera su vida útil estimada.",
            "capex_lines": [{"time_horizon_code": "MEDIO", "amount": "83407.50"}],
        },
    ).json()
    cliente.post(
        f"{RUTA}/findings/{hallazgo['id']}/transitions",
        headers=cab("consultor_a"),
        json={"to": "EN_REVISION"},
    )
    return {"asset": activo, "finding": hallazgo}


def generar(
    cliente: TestClient, cab: Any, proyecto: str, plantilla: dict[str, Any], **extra: Any
) -> Any:
    return cliente.post(
        f"{RUTA}/projects/{proyecto}/reports",
        headers=cab("admin_a"),
        json={"template_id": plantilla["id"], **extra},
    )


# ─────────────────────────────────────────────────────────────────────────────
#  Plantillas
# ─────────────────────────────────────────────────────────────────────────────


def test_la_plantilla_se_analiza_al_registrarla(plantilla: dict[str, Any]) -> None:
    """Analizar al registrar y no al generar: si la plantilla tiene problemas,
    se descubren ahora y no cuando alguien espera un informe."""
    analisis = plantilla["analysis"]
    assert analisis["slide_count"] >= 1
    assert set(analisis["placeholders"]) >= {"project.code", "project.name", "capex.total"}


def test_el_analisis_avisa_de_la_marca_de_agua(plantilla: dict[str, Any]) -> None:
    """`[REQ]` P-43 · La generación la retirará, y conviene que nadie se
    sorprenda al comparar con la plantilla."""
    assert plantilla["analysis"]["has_watermark"] is True


def test_un_fichero_que_no_es_pptx_se_rechaza(cliente: TestClient, cab: Any) -> None:
    r = cliente.post(
        f"{RUTA}/report-templates",
        headers=cab("admin_a"),
        files={
            "file": ("falso.pptx", io.BytesIO(b"no soy un zip"), "application/vnd.ms-powerpoint")
        },
        data={"name": "Falsa"},
    )
    assert r.status_code == 415


# ─────────────────────────────────────────────────────────────────────────────
#  Mapeo
# ─────────────────────────────────────────────────────────────────────────────


def test_una_expresion_inexistente_se_rechaza_al_guardar(
    cliente: TestClient, cab: Any, plantilla: dict[str, Any]
) -> None:
    """Descubrirlo al generar es tarde: §17.7 lo declara bloqueante justo por
    eso, y aquí se evita que llegue a existir."""
    r = cliente.post(
        f"{RUTA}/report-templates/{plantilla['id']}/mappings",
        headers=cab("admin_a"),
        json={"name": "Roto", "bindings": {"project.code": "project.inventado"}},
    )
    assert r.status_code == 422
    assert "project.inventado" in r.json()["detail"]


def test_solo_hay_un_mapeo_por_defecto(
    cliente: TestClient, cab: Any, plantilla: dict[str, Any], mapeo: dict[str, Any]
) -> None:
    cliente.post(
        f"{RUTA}/report-templates/{plantilla['id']}/mappings",
        headers=cab("admin_a"),
        json={"name": "Otro", "bindings": MAPEO_COMPLETO, "is_default": True},
    )
    mapeos = cliente.get(
        f"{RUTA}/report-templates/{plantilla['id']}/mappings", headers=cab("admin_a")
    ).json()
    assert sum(1 for m in mapeos if m["is_default"]) == 1


# ─────────────────────────────────────────────────────────────────────────────
#  Avisos previos
# ─────────────────────────────────────────────────────────────────────────────


def test_sin_mapeo_los_marcadores_bloquean(
    cliente: TestClient, cab: Any, proyecto: str, plantilla: dict[str, Any]
) -> None:
    r = cliente.post(
        f"{RUTA}/projects/{proyecto}/reports/preflight",
        headers=cab("admin_a"),
        json={"template_id": plantilla["id"]},
    )
    assert r.status_code == 200
    assert r.json()["can_generate"] is False
    assert any(w["codigo"] == "UNMAPPED_PLACEHOLDER" for w in r.json()["warnings"])


def test_con_el_mapeo_completo_ya_se_puede_generar(
    cliente: TestClient, cab: Any, proyecto: str, plantilla: dict[str, Any], mapeo: dict[str, Any]
) -> None:
    r = cliente.post(
        f"{RUTA}/projects/{proyecto}/reports/preflight",
        headers=cab("admin_a"),
        json={"template_id": plantilla["id"], "mapping_id": mapeo["id"]},
    )
    assert r.json()["can_generate"] is True


def test_los_precios_sin_validar_avisan_sin_impedir_generar(
    cliente: TestClient,
    cab: Any,
    proyecto: str,
    plantilla: dict[str, Any],
    mapeo: dict[str, Any],
    con_hallazgo: dict[str, Any],
) -> None:
    """`[REC]` Un borrador interno con precios sin validar es legítimo; lo que
    no puede pasar es enviarlo al cliente sin darse cuenta."""
    r = cliente.post(
        f"{RUTA}/projects/{proyecto}/reports/preflight",
        headers=cab("admin_a"),
        json={"template_id": plantilla["id"], "mapping_id": mapeo["id"]},
    ).json()
    assert r["can_generate"] is True
    aviso = next(w for w in r["warnings"] if w["codigo"] == "UNVALIDATED_PRICES")
    assert aviso["bloquea"] is False
    # En español, no «83,407.50»: la aplicación entera va en español.
    assert "83.407,50 €" in aviso["mensaje"]


def test_un_hallazgo_en_borrador_avisa_de_que_no_saldra_en_el_informe(
    cliente: TestClient,
    cab: Any,
    proyecto: str,
    plantilla: dict[str, Any],
    mapeo: dict[str, Any],
    motor_admin: Engine,
) -> None:
    """El snapshot solo publica `EN_REVISION` y `VALIDADO`. Sin este aviso, un
    encargo con todo en borrador generaba un informe con la tabla de CAPEX
    vacía y «CAPEX total: 0,00 €», sin advertir de nada."""
    with motor_admin.begin() as conn:
        tipologia = conn.execute(
            text("SELECT id FROM asset_typology ORDER BY sort_order LIMIT 1")
        ).scalar_one()
        zona = conn.execute(
            text(
                "SELECT z.id FROM zone z JOIN zone_typology zt ON zt.zone_id = z.id "
                "WHERE zt.typology_id = :t LIMIT 1"
            ),
            {"t": tipologia},
        ).scalar_one()
        codigo = conn.execute(
            text("SELECT id FROM capex_code WHERE level = 3 LIMIT 1")
        ).scalar_one()

    activo = cliente.post(
        f"{RUTA}/projects/{proyecto}/assets",
        headers=cab("consultor_a"),
        json={"name": "Nave en borrador", "typology_id": str(tipologia)},
    ).json()
    # Se queda en BORRADOR: no se pide ninguna transición.
    cliente.post(
        f"{RUTA}/projects/{proyecto}/findings",
        headers=cab("consultor_a"),
        json={
            "asset_id": activo["id"],
            "capex_code_id": str(codigo),
            "zone_id": str(zona),
            "title": "Sustitución de enfriadora",
            "description": "Equipo fuera de reglamento.",
            "capex_lines": [{"time_horizon_code": "CORTO", "amount": "271700.00"}],
        },
    )

    r = cliente.post(
        f"{RUTA}/projects/{proyecto}/reports/preflight",
        headers=cab("admin_a"),
        json={"template_id": plantilla["id"], "mapping_id": mapeo["id"]},
    ).json()
    aviso = next(w for w in r["warnings"] if w["codigo"] == "DRAFT_FINDINGS_EXCLUDED")
    # No bloquea: un Red Flag temprano con todo en borrador es legítimo.
    assert aviso["bloquea"] is False
    assert r["can_generate"] is True
    assert "1 hallazgos" in aviso["mensaje"]
    assert "271.700,00 €" in aviso["mensaje"]


def test_un_hallazgo_ya_en_revision_no_dispara_el_aviso_de_borradores(
    cliente: TestClient,
    cab: Any,
    proyecto: str,
    plantilla: dict[str, Any],
    mapeo: dict[str, Any],
    con_hallazgo: dict[str, Any],
) -> None:
    """`con_hallazgo` deja el hallazgo en `EN_REVISION`: sí sale en el informe,
    así que no hay nada que declarar."""
    r = cliente.post(
        f"{RUTA}/projects/{proyecto}/reports/preflight",
        headers=cab("admin_a"),
        json={"template_id": plantilla["id"], "mapping_id": mapeo["id"]},
    ).json()
    assert not any(w["codigo"] == "DRAFT_FINDINGS_EXCLUDED" for w in r["warnings"])


# ─────────────────────────────────────────────────────────────────────────────
#  Generación
# ─────────────────────────────────────────────────────────────────────────────


def test_se_genera_el_informe_con_su_pptx_y_su_xlsx(
    cliente: TestClient,
    cab: Any,
    proyecto: str,
    plantilla: dict[str, Any],
    mapeo: dict[str, Any],
    con_hallazgo: dict[str, Any],
) -> None:
    r = generar(cliente, cab, proyecto, plantilla, mapping_id=mapeo["id"])
    assert r.status_code == 201, r.text
    version = r.json()
    assert version["version_number"] == 1
    assert version["status"] == "GENERADO"
    assert len(version["pptx_sha256"]) == 64
    assert len(version["data_snapshot_sha256"]) == 64

    for formato in ("pptx", "xlsx"):
        descarga = cliente.get(
            f"{RUTA}/reports/{version['id']}/download?formato={formato}", headers=cab("admin_a")
        )
        assert descarga.status_code == 200
        assert descarga.content[:2] == b"PK", "es un contenedor OOXML"


def test_los_marcadores_se_sustituyen_de_verdad(
    cliente: TestClient,
    cab: Any,
    proyecto: str,
    plantilla: dict[str, Any],
    mapeo: dict[str, Any],
    con_hallazgo: dict[str, Any],
) -> None:
    """Se comprueba abriendo el PPTX generado, no confiando en el contador."""
    version = generar(cliente, cab, proyecto, plantilla, mapping_id=mapeo["id"]).json()
    datos = cliente.get(f"{RUTA}/reports/{version['id']}/download", headers=cab("admin_a")).content

    prs = Presentation(io.BytesIO(datos))
    texto = "\n".join(
        forma.text_frame.text
        for slide in prs.slides
        for forma in slide.shapes
        if forma.has_text_frame
    )
    assert "{{" not in texto, "no queda ningún marcador sin resolver"
    assert "Encargo con informe" in texto
    assert "Inversora Ficticia" in texto


def test_la_marca_de_agua_no_aparece_en_lo_generado(
    cliente: TestClient,
    cab: Any,
    proyecto: str,
    plantilla: dict[str, Any],
    mapeo: dict[str, Any],
) -> None:
    """`[REQ]` P-43 · «La marca de agua de DRAFT no deberá aparecer en las
    versiones futuras, retírala.»"""
    from tdd.reporting.watermark import hay_marcas_de_agua

    version = generar(cliente, cab, proyecto, plantilla, mapping_id=mapeo["id"]).json()
    datos = cliente.get(f"{RUTA}/reports/{version['id']}/download", headers=cab("admin_a")).content
    assert hay_marcas_de_agua(Presentation(io.BytesIO(datos))) is False


def test_la_plantilla_original_no_se_toca(
    cliente: TestClient,
    cab: Any,
    proyecto: str,
    plantilla: dict[str, Any],
    mapeo: dict[str, Any],
    motor_admin: Engine,
) -> None:
    """`[REQ]` Se comprueba por hash, antes y después de generar."""
    antes = plantilla["sha256"]
    generar(cliente, cab, proyecto, plantilla, mapping_id=mapeo["id"])
    with motor_admin.begin() as conn:
        despues = conn.execute(
            text("SELECT sha256 FROM report_template WHERE id = :i"), {"i": plantilla["id"]}
        ).scalar_one()
    assert despues == antes


def test_las_fotos_seleccionadas_entran_en_el_informe(
    cliente: TestClient,
    cab: Any,
    proyecto: str,
    plantilla: dict[str, Any],
    mapeo: dict[str, Any],
    con_hallazgo: dict[str, Any],
) -> None:
    foto = cliente.post(
        f"{RUTA}/projects/{proyecto}/photos",
        headers=cab("consultor_a"),
        files={"file": ("c.jpg", io.BytesIO(imagen(color=(3, 33, 133))), "image/jpeg")},
        data={"asset_id": con_hallazgo["asset"]["id"]},
    ).json()
    cliente.patch(
        f"{RUTA}/photos/{foto['id']}",
        headers=cab("consultor_a"),
        json={"include_in_report": True, "caption": "Cubierta invertida", "report_order": 1},
    )

    version = generar(cliente, cab, proyecto, plantilla, mapping_id=mapeo["id"]).json()
    datos = cliente.get(f"{RUTA}/reports/{version['id']}/download", headers=cab("admin_a")).content

    prs = Presentation(io.BytesIO(datos))
    imagenes = [f for slide in prs.slides for f in slide.shapes if f.shape_type == 13]
    assert imagenes, "la fotografía se ha insertado"

    pies = "\n".join(
        f.text_frame.text for slide in prs.slides for f in slide.shapes if f.has_text_frame
    )
    assert "Cubierta invertida" in pies


def test_las_anotaciones_llegan_quemadas_al_informe(
    cliente: TestClient,
    cab: Any,
    proyecto: str,
    plantilla: dict[str, Any],
    mapeo: dict[str, Any],
    con_hallazgo: dict[str, Any],
) -> None:
    """`[REQ]` §15.2 · Señalar la fisura con una flecha es lo que hace útil una
    foto técnica.

    La capa se guardaba desde el principio —versionada, auditada, reversible—
    pero **el generador insertaba la foto sin ella**: anotar producía un JSON
    que no llegaba a ninguna parte.

    Se comprueba generando dos veces la **misma** foto, antes y después de
    anotarla, y mirando si el binario que acaba dentro del PPTX cambia. Contar
    píxeles de un color no vale: la imagen de prueba es un degradado y la
    plantilla trae sus propios logotipos, así que siempre hay «rojo» en alguna
    parte y la comprobación pasaría sin arreglar nada. Se verificó desactivando
    el rasterizado a propósito: así la prueba falla, como debe.
    """
    foto = cliente.post(
        f"{RUTA}/projects/{proyecto}/photos",
        headers=cab("consultor_a"),
        files={"file": ("a.jpg", io.BytesIO(imagen(color=(90, 140, 190))), "image/jpeg")},
        data={"asset_id": con_hallazgo["asset"]["id"]},
    ).json()
    cliente.patch(
        f"{RUTA}/photos/{foto['id']}",
        headers=cab("consultor_a"),
        json={"include_in_report": True, "report_order": 1},
    )

    def imagenes_del_informe() -> set[bytes]:
        version = generar(cliente, cab, proyecto, plantilla, mapping_id=mapeo["id"]).json()
        datos = cliente.get(
            f"{RUTA}/reports/{version['id']}/download", headers=cab("admin_a")
        ).content
        prs = Presentation(io.BytesIO(datos))
        return {
            f.image.blob
            for slide in prs.slides
            for f in slide.shapes
            if f.shape_type == 13  # noqa: PLR2004 — PICTURE
        }

    sin_anotar = imagenes_del_informe()
    assert sin_anotar, "la fotografía se ha insertado"

    r = cliente.post(
        f"{RUTA}/photos/{foto['id']}/versions/annotate",
        headers=cab("consultor_a"),
        json={
            "annotations": {
                "shapes": [
                    {
                        "tipo": "FLECHA",
                        "x1": 0.1,
                        "y1": 0.1,
                        "x2": 0.8,
                        "y2": 0.8,
                        "color": "#DC2626",
                        "grosor": 8,
                    }
                ]
            }
        },
    )
    assert r.status_code == 201, r.text

    con_anotar = imagenes_del_informe()
    nuevas = con_anotar - sin_anotar
    assert nuevas, "el PPTX lleva exactamente la misma imagen: la anotación no se ha pintado"


def test_una_foto_en_cuarentena_impide_generar(
    cliente: TestClient,
    cab: Any,
    proyecto: str,
    plantilla: dict[str, Any],
    mapeo: dict[str, Any],
    motor_admin: Engine,
) -> None:
    """Es uno de los cinco bloqueantes: produciría un documento con evidencia
    no verificada."""
    foto = cliente.post(
        f"{RUTA}/projects/{proyecto}/photos",
        headers=cab("consultor_a"),
        files={"file": ("q.jpg", io.BytesIO(imagen(color=(9, 9, 99))), "image/jpeg")},
    ).json()
    cliente.patch(
        f"{RUTA}/photos/{foto['id']}", headers=cab("consultor_a"), json={"include_in_report": True}
    )
    with motor_admin.begin() as conn:
        conn.execute(
            text("UPDATE photo SET status = 'CUARENTENA' WHERE id = :i"), {"i": foto["id"]}
        )

    r = generar(cliente, cab, proyecto, plantilla, mapping_id=mapeo["id"])
    assert r.status_code == 409
    assert "verificaciones" in r.json()["detail"]


# ─────────────────────────────────────────────────────────────────────────────
#  Snapshot §17.6
# ─────────────────────────────────────────────────────────────────────────────


def test_el_informe_sigue_diciendo_lo_mismo_aunque_los_datos_cambien(
    cliente: TestClient,
    cab: Any,
    proyecto: str,
    plantilla: dict[str, Any],
    mapeo: dict[str, Any],
    con_hallazgo: dict[str, Any],
    motor_admin: Engine,
) -> None:
    """`[REQ]` §17.6 · **La garantía que estructura el bloque entero.**

    La generación lee del snapshot, no de la base. Un cambio posterior —o
    concurrente— no puede alterar lo que ya se entregó.
    """
    version = generar(cliente, cab, proyecto, plantilla, mapping_id=mapeo["id"]).json()
    with motor_admin.begin() as conn:
        antes = conn.execute(
            text("SELECT data_snapshot FROM report_version WHERE id = :i"), {"i": version["id"]}
        ).scalar_one()

    # Los datos siguen cambiando después de generar.
    linea = con_hallazgo["finding"]["capex_lines"][0]
    cliente.patch(
        f"{RUTA}/capex-items/{linea['id']}", headers=cab("consultor_a"), json={"amount": "1.00"}
    )

    with motor_admin.begin() as conn:
        despues = conn.execute(
            text("SELECT data_snapshot FROM report_version WHERE id = :i"), {"i": version["id"]}
        ).scalar_one()
    assert despues == antes
    assert despues["capex_items"][0]["amount"] == "83407.5000"


def test_el_snapshot_congela_los_catalogos(
    cliente: TestClient,
    cab: Any,
    proyecto: str,
    plantilla: dict[str, Any],
    mapeo: dict[str, Any],
    con_hallazgo: dict[str, Any],
    motor_admin: Engine,
) -> None:
    """`[REC]` Sin esto, retirar un código CAPEX dentro de dos años dejaría
    huecos en un informe ya entregado."""
    version = generar(cliente, cab, proyecto, plantilla, mapping_id=mapeo["id"]).json()
    with motor_admin.begin() as conn:
        congelado = conn.execute(
            text("SELECT data_snapshot FROM report_version WHERE id = :i"), {"i": version["id"]}
        ).scalar_one()

    catalogos = congelado["catalogs"]
    assert catalogos["risk_levels"], "los grados de riesgo, con su definición íntegra"
    assert all(r["definition_es"] for r in catalogos["risk_levels"])
    assert catalogos["capex_codes"], "los códigos realmente usados"


def test_un_borrador_no_sale_en_el_informe(
    cliente: TestClient,
    cab: Any,
    proyecto: str,
    plantilla: dict[str, Any],
    mapeo: dict[str, Any],
    con_hallazgo: dict[str, Any],
    motor_admin: Engine,
) -> None:
    """Lo que aún se está escribiendo no debe aparecer en un documento que se
    entrega al cliente."""
    cliente.post(
        f"{RUTA}/findings/{con_hallazgo['finding']['id']}/transitions",
        headers=cab("consultor_a"),
        json={"to": "BORRADOR"},
    )
    version = generar(cliente, cab, proyecto, plantilla, mapping_id=mapeo["id"]).json()
    with motor_admin.begin() as conn:
        congelado = conn.execute(
            text("SELECT data_snapshot FROM report_version WHERE id = :i"), {"i": version["id"]}
        ).scalar_one()
    assert congelado["findings"] == []


def test_la_segunda_version_sustituye_a_la_primera(
    cliente: TestClient,
    cab: Any,
    proyecto: str,
    plantilla: dict[str, Any],
    mapeo: dict[str, Any],
    con_hallazgo: dict[str, Any],
) -> None:
    v1 = generar(cliente, cab, proyecto, plantilla, mapping_id=mapeo["id"]).json()
    v2 = generar(cliente, cab, proyecto, plantilla, mapping_id=mapeo["id"]).json()
    assert v2["version_number"] == 2
    assert v2["supersedes_version_id"] == v1["id"]


def test_la_comparacion_dice_cuanto_se_movio_el_capex(
    cliente: TestClient,
    cab: Any,
    proyecto: str,
    plantilla: dict[str, Any],
    mapeo: dict[str, Any],
    con_hallazgo: dict[str, Any],
) -> None:
    """`[REQ]` §17.6 · Lo que de verdad se mira al comparar dos versiones."""
    v1 = generar(cliente, cab, proyecto, plantilla, mapping_id=mapeo["id"]).json()
    linea = con_hallazgo["finding"]["capex_lines"][0]
    cliente.patch(
        f"{RUTA}/capex-items/{linea['id']}", headers=cab("consultor_a"), json={"amount": "90000.00"}
    )
    v2 = generar(cliente, cab, proyecto, plantilla, mapping_id=mapeo["id"]).json()

    diff = cliente.get(f"{RUTA}/reports/{v1['id']}/diff/{v2['id']}", headers=cab("admin_a")).json()
    medio = diff["capex_by_horizon"]["MEDIO"]
    assert medio["before"] == "83407.5000"
    assert medio["after"] == "90000.0000"
    assert medio["delta"] == "6592.5000"


# ─────────────────────────────────────────────────────────────────────────────
#  Ciclo de vida e inmutabilidad §17.6
# ─────────────────────────────────────────────────────────────────────────────


def test_el_informe_recorre_su_ciclo_hasta_emitirse(
    cliente: TestClient,
    cab: Any,
    proyecto: str,
    plantilla: dict[str, Any],
    mapeo: dict[str, Any],
    con_hallazgo: dict[str, Any],
) -> None:
    version = generar(cliente, cab, proyecto, plantilla, mapping_id=mapeo["id"]).json()
    for destino in ("EN_REVISION", "APROBADO", "EMITIDO"):
        r = cliente.post(
            f"{RUTA}/reports/{version['id']}/transitions",
            headers=cab("admin_a"),
            json={"to": destino},
        )
        assert r.status_code == 200, r.text
        assert r.json()["status"] == destino
    assert r.json()["is_locked"] is True


def test_emitir_saltandose_la_revision_es_un_conflicto(
    cliente: TestClient,
    cab: Any,
    proyecto: str,
    plantilla: dict[str, Any],
    mapeo: dict[str, Any],
    con_hallazgo: dict[str, Any],
) -> None:
    version = generar(cliente, cab, proyecto, plantilla, mapping_id=mapeo["id"]).json()
    r = cliente.post(
        f"{RUTA}/reports/{version['id']}/transitions",
        headers=cab("admin_a"),
        json={"to": "EMITIDO"},
    )
    assert r.status_code == 409


def test_un_consultor_no_puede_aprobar_ni_emitir(
    cliente: TestClient,
    cab: Any,
    proyecto: str,
    plantilla: dict[str, Any],
    mapeo: dict[str, Any],
    con_hallazgo: dict[str, Any],
) -> None:
    version = generar(cliente, cab, proyecto, plantilla, mapping_id=mapeo["id"]).json()
    cliente.post(
        f"{RUTA}/reports/{version['id']}/transitions",
        headers=cab("admin_a"),
        json={"to": "EN_REVISION"},
    )
    r = cliente.post(
        f"{RUTA}/reports/{version['id']}/transitions",
        headers=cab("consultor_a"),
        json={"to": "APROBADO"},
    )
    assert r.status_code == 403


def test_un_informe_emitido_no_se_puede_cambiar_por_la_api(
    cliente: TestClient,
    cab: Any,
    proyecto: str,
    plantilla: dict[str, Any],
    mapeo: dict[str, Any],
    con_hallazgo: dict[str, Any],
) -> None:
    version = generar(cliente, cab, proyecto, plantilla, mapping_id=mapeo["id"]).json()
    for destino in ("EN_REVISION", "APROBADO", "EMITIDO"):
        cliente.post(
            f"{RUTA}/reports/{version['id']}/transitions",
            headers=cab("admin_a"),
            json={"to": destino},
        )
    r = cliente.post(
        f"{RUTA}/reports/{version['id']}/transitions",
        headers=cab("admin_a"),
        json={"to": "EN_REVISION"},
    )
    assert r.status_code == 409
    assert "inmutable" in r.json()["detail"]


def test_la_base_de_datos_tampoco_deja_tocarlo(
    cliente: TestClient,
    cab: Any,
    proyecto: str,
    plantilla: dict[str, Any],
    mapeo: dict[str, Any],
    con_hallazgo: dict[str, Any],
    motor_admin: Engine,
) -> None:
    """La garantía de verdad. Si alguien escribe un `UPDATE` nuevo dentro de
    seis meses, esto sigue en pie."""
    version = generar(cliente, cab, proyecto, plantilla, mapping_id=mapeo["id"]).json()
    for destino in ("EN_REVISION", "APROBADO", "EMITIDO"):
        cliente.post(
            f"{RUTA}/reports/{version['id']}/transitions",
            headers=cab("admin_a"),
            json={"to": destino},
        )
    with motor_admin.begin() as conn, pytest.raises(Exception, match="inmutable"):
        conn.execute(
            text("UPDATE report_version SET status = 'GENERADO' WHERE id = :i"),
            {"i": version["id"]},
        )


def test_un_informe_emitido_no_se_borra(
    cliente: TestClient,
    cab: Any,
    proyecto: str,
    plantilla: dict[str, Any],
    mapeo: dict[str, Any],
    con_hallazgo: dict[str, Any],
    motor_admin: Engine,
) -> None:
    version = generar(cliente, cab, proyecto, plantilla, mapping_id=mapeo["id"]).json()
    for destino in ("EN_REVISION", "APROBADO", "EMITIDO"):
        cliente.post(
            f"{RUTA}/reports/{version['id']}/transitions",
            headers=cab("admin_a"),
            json={"to": destino},
        )
    with motor_admin.begin() as conn, pytest.raises(Exception, match="no se borra"):
        conn.execute(text("DELETE FROM report_version WHERE id = :i"), {"i": version["id"]})


def test_el_informe_emitido_sigue_descargandose_igual(
    cliente: TestClient,
    cab: Any,
    proyecto: str,
    plantilla: dict[str, Any],
    mapeo: dict[str, Any],
    con_hallazgo: dict[str, Any],
) -> None:
    """Es lo que hace verificable el `pptx_sha256` que se guardó al generarlo."""
    from tdd.evidence.images import sha256_de

    version = generar(cliente, cab, proyecto, plantilla, mapping_id=mapeo["id"]).json()
    for destino in ("EN_REVISION", "APROBADO", "EMITIDO"):
        cliente.post(
            f"{RUTA}/reports/{version['id']}/transitions",
            headers=cab("admin_a"),
            json={"to": destino},
        )
    datos = cliente.get(f"{RUTA}/reports/{version['id']}/download", headers=cab("admin_a")).content
    assert sha256_de(datos) == version["pptx_sha256"]


def test_otra_organizacion_no_ve_el_informe(
    cliente: TestClient,
    cab: Any,
    proyecto: str,
    plantilla: dict[str, Any],
    mapeo: dict[str, Any],
    con_hallazgo: dict[str, Any],
) -> None:
    version = generar(cliente, cab, proyecto, plantilla, mapping_id=mapeo["id"]).json()
    assert cliente.get(f"{RUTA}/reports/{version['id']}", headers=cab("admin_b")).status_code == 404


def test_la_generacion_queda_auditada(
    cliente: TestClient,
    cab: Any,
    proyecto: str,
    plantilla: dict[str, Any],
    mapeo: dict[str, Any],
    con_hallazgo: dict[str, Any],
    motor_admin: Engine,
) -> None:
    version = generar(cliente, cab, proyecto, plantilla, mapping_id=mapeo["id"]).json()
    with motor_admin.begin() as conn:
        acciones = (
            conn.execute(
                text("SELECT action FROM audit_log WHERE entity_id = :i"), {"i": version["id"]}
            )
            .scalars()
            .all()
        )
    assert "REPORT_GENERATED" in acciones


def test_el_pptx_y_el_excel_del_informe_llevan_las_mismas_actuaciones(
    cliente: TestClient,
    cab: Any,
    proyecto: str,
    plantilla: dict[str, Any],
    mapeo: dict[str, Any],
    con_hallazgo: dict[str, Any],
) -> None:
    """`[REQ]` P-31 · Los dos ficheros viajan en el mismo correo.

    Antes esto se garantizaba compartiendo `CapexTableLayout`, pero el Excel
    pasó a ser **la plantilla del cliente rellenada** y sus columnas son otras
    a propósito. Lo que sigue teniendo que cuadrar es el contenido: una
    actuación que sale en la diapositiva y no en la hoja —o al revés— es dinero
    que aparece o desaparece según qué fichero abra quien lo recibe.
    """
    import io
    import zipfile

    from openpyxl import load_workbook
    from pptx import Presentation

    version = generar(cliente, cab, proyecto, plantilla, mapping_id=mapeo["id"]).json()
    pptx = cliente.get(f"{RUTA}/reports/{version['id']}/download", headers=cab("admin_a")).content
    xlsx = cliente.get(
        f"{RUTA}/reports/{version['id']}/download?formato=xlsx", headers=cab("admin_a")
    ).content

    titulo = "Renovación de impermeabilización"

    en_pptx = any(
        titulo in celda.text
        for slide in Presentation(io.BytesIO(pptx)).slides
        for forma in slide.shapes
        if forma.has_table
        for fila in forma.table.rows
        for celda in fila.cells
    )
    hoja = load_workbook(io.BytesIO(xlsx))["CapEx"]
    en_xlsx = any(hoja[f"G{f}"].value == titulo for f in range(14, 281))

    assert en_pptx, "la actuación no sale en la tabla del PowerPoint"
    assert en_xlsx, "la actuación no sale en la hoja del Excel"
    # Y el Excel es la plantilla del cliente, no un libro construido a mano.
    with zipfile.ZipFile(io.BytesIO(xlsx)) as z:
        assert "xl/pivotTables/pivotTable1.xml" in z.namelist()
