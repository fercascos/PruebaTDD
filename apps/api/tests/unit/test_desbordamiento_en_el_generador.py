"""Que el aviso de desbordamiento **llegue al informe**, no solo al módulo.

`overflow.py` estaba construido, medido con la Gotham real y probado en
`test_fuentes_y_desbordamiento.py`… y **no lo llamaba nadie**. El generador no
lo usaba, así que el aviso no se emitía jamás sobre un informe de verdad: el
consultor se enteraba de que un texto no cabía abriendo el PPTX, o el cliente.

Estas pruebas fijan el enganche. Se construye una plantilla mínima en memoria
—no hace falta la del cliente para esto— con un marco pequeño y un marcador
dentro, y se comprueba que un texto largo produce aviso y uno corto no.
"""

from __future__ import annotations

import io

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from tdd.reporting import fonts
from tdd.reporting.clone import sustituir_marcadores
from tdd.reporting.generator import _evaluar_marco

#: Sin las corporativas instaladas no hay medición, y el módulo prefiere
#: callarse antes que medir con una sustituta y dar un número inventado. Estas
#: pruebas comprueban justamente que se mide, así que sin fuente no aplican.
sin_gotham = pytest.mark.skipif(
    fonts.localizar("Gotham Light") is None,
    reason="Gotham no instalada: sin ella no se mide, y eso ya lo cubre otra prueba",
)


def _plantilla(ancho_in: float, alto_in: float) -> bytes:
    """Una diapositiva con un solo marco, del tamaño que se pida."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # en blanco
    caja = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(ancho_in), Inches(alto_in))
    run = caja.text_frame.paragraphs[0].add_run()
    run.text = "{{texto}}"
    run.font.name = "Gotham Light"
    run.font.size = Pt(10)
    salida = io.BytesIO()
    prs.save(salida)
    return salida.getvalue()


def _avisos(plantilla: bytes, texto: str) -> list[str]:
    prs = Presentation(io.BytesIO(plantilla))
    recogidos: list[str] = []

    def medir(forma: object) -> None:
        aviso = _evaluar_marco(forma)
        if aviso is not None:
            recogidos.append(aviso)

    for slide in prs.slides:
        sustituir_marcadores(slide, {"texto": texto}, medir=medir)
    return recogidos


@sin_gotham
def test_un_texto_que_no_cabe_produce_aviso() -> None:
    """`[REQ]` Es para lo que existe toda la medición con la fuente real."""
    # Un marco de 2×1 pulgadas a 10 pt cabe del orden de 200 caracteres.
    avisos = _avisos(_plantilla(2, 1), "palabra " * 400)
    assert avisos, "El aviso no llegó al informe: el enganche no funciona"
    assert "no cabe" in avisos[0]
    # Y dice **cuánto** sobra, que es lo accionable: «recorta unos N».
    assert "sobran" in avisos[0]


@sin_gotham
def test_un_texto_que_cabe_no_molesta() -> None:
    """Un aviso que salta siempre es un aviso que se acaba ignorando."""
    assert _avisos(_plantilla(8, 5), "Cabe de sobra.") == []


@sin_gotham
def test_avisa_tambien_cuando_va_justo() -> None:
    """`CERCA` existe porque el margen de la estimación es de un ±10 %.

    Un texto al 95 % del marco puede caber o no según cómo parta las líneas
    PowerPoint, y eso es exactamente lo que hay que revisar a mano antes de
    entregar.
    """
    # 2×1 in a 10 pt ≈ 200 caracteres: se apunta a ocupar poco más del 90 %.
    avisos = _avisos(_plantilla(2, 1), "x" * 190)
    assert avisos and ("%" in avisos[0])


def test_sin_la_fuente_declarada_no_se_inventa_un_numero() -> None:
    """`[REQ]` Medir con una sustituta es peor que no medir.

    Un marco cuya tipografía no está instalada no produce aviso: la capacidad
    saldría de una heurística y presentarla como una medición sería dar un
    número inventado con aspecto de dato.
    """
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    caja = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(2), Inches(1))
    run = caja.text_frame.paragraphs[0].add_run()
    run.text = "{{texto}}"
    run.font.name = "Tipografía Que No Existe En Ningún Sitio"
    run.font.size = Pt(10)
    salida = io.BytesIO()
    prs.save(salida)

    assert _avisos(salida.getvalue(), "palabra " * 400) == []
