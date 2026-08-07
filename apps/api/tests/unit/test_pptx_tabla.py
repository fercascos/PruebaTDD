"""La tabla de CAPEX dibujada en la diapositiva.

`capex_layout` decide **qué** filas y columnas hay; este módulo decide cómo se
ven. Hasta ahora no tenía prueba ninguna, y se cambió al pasar la tabla al
desglose de la plantilla CAPEX vigente: es justo cuando hace falta.

La prueba que sostiene el fichero es `test_la_tabla_cabe_en_la_diapositiva`.
La plantilla del cliente es **4:3 (10 × 7,5 in)** y ahí está la restricción que
manda sobre todo lo demás: es la razón de que el tipo de coste y el capítulo
sean filas de sección en vez de columnas.
"""

from __future__ import annotations

from decimal import Decimal

import pytest
from pptx import Presentation
from pptx.util import Inches

import tdd.reporting.capex_layout as cl
from tdd.reporting.pptx_table import GRIS_SECCION, GRIS_TIPO_DE_COSTE, insertar_tabla

#: La plantilla del cliente. Ver `docs/18` §«Tamaño».
ANCHO_DIAPOSITIVA_IN = 10.0
ALTO_DIAPOSITIVA_IN = 7.5


def _linea(**kw) -> cl.LineaCapex:
    base = dict(
        numero="",
        zona="Cubierta",
        concepto="Vida útil",
        descripcion="Cubierta deck al final de su vida útil",
        riesgo="04 Extremo",
        comentarios="Lámina asfáltica agotada",
        horizonte="CORTO",
        importe=Decimal("412500"),
        tipo_de_coste="Hard Costs",
        capitulo="Cubierta",
        objeto="Cubierta",
        recuperable="NO",
    )
    return cl.LineaCapex(**{**base, **kw})


LINEAS = [
    _linea(),
    _linea(
        capitulo="HVAC",
        objeto="Producción de climatización",
        zona="Cuartos técnicos",
        descripcion="Enfriadora fuera de reglamento",
        importe=Decimal("271700"),
        concepto="Normativa",
    ),
    _linea(
        tipo_de_coste="Soft Costs",
        capitulo="Licencias y Tasas",
        objeto="General",
        zona="General",
        descripcion="Licencia de obras",
        riesgo="-",
        horizonte="OTRO",
        importe=Decimal("18400"),
        concepto="Soft Cost",
        recuperable="N.A.",
    ),
]


@pytest.fixture
def tabla():
    layout = cl.construir(LINEAS, capitulo="Instalaciones")
    prs = Presentation()
    prs.slide_width = Inches(ANCHO_DIAPOSITIVA_IN)
    prs.slide_height = Inches(ALTO_DIAPOSITIVA_IN)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    insertar_tabla(slide, layout, left_in=0.3, top_in=0.6)
    return slide.shapes[0]


def _texto(t, fila: int) -> list[str]:
    return [t.cell(fila, i).text for i in range(len(t.columns))]


# ─────────────────────────────────────────────────────────────────────────────
#  La restricción que manda: 4:3
# ─────────────────────────────────────────────────────────────────────────────


def test_la_tabla_cabe_en_la_diapositiva(tabla) -> None:
    """`[LIM]` 10 in de ancho, y la tabla se coloca a 0,3 in del borde. Si esto
    falla, el informe sale con columnas fuera de la página: no es un detalle
    estético, es contenido que el cliente no ve."""
    ancho_in = tabla.width / 914400
    assert ancho_in + 0.3 <= ANCHO_DIAPOSITIVA_IN


def test_las_columnas_suman_el_ancho_de_la_tabla(tabla) -> None:
    suma = sum(c.width for c in tabla.table.columns) / 914400
    assert round(suma, 2) == round(tabla.width / 914400, 2)


# ─────────────────────────────────────────────────────────────────────────────
#  Cabecera de dos niveles
# ─────────────────────────────────────────────────────────────────────────────


def test_los_cinco_plazos_cuelgan_de_una_cabecera_combinada(tabla) -> None:
    fila1 = _texto(tabla.table, 1)
    assert "CAPEX ESTIMADO" in fila1
    # Los subtítulos van debajo, en la segunda fila de cabecera.
    fila2 = _texto(tabla.table, 2)
    assert [x for x in fila2 if x] == [
        "Corto plazo",
        "Medio plazo",
        "Largo plazo",
        "Mejoras",
        "Otro",
    ]


def test_el_total_no_lo_absorbe_la_cabecera_de_plazos(tabla) -> None:
    """`TOTAL` va detrás de los cinco plazos pero **no cuelga** de ellos: es la
    suma, no un plazo más. Si la combinación se lo tragase, la columna se
    quedaría sin título."""
    assert _texto(tabla.table, 1)[-1] == "TOTAL"


# ─────────────────────────────────────────────────────────────────────────────
#  Los dos niveles de sección
# ─────────────────────────────────────────────────────────────────────────────


def test_el_tipo_de_coste_y_el_capitulo_se_distinguen_por_tono(tabla) -> None:
    """Con el mismo gris, una tabla de cinco capítulos parecería tener diez
    secciones sueltas en vez de dos grupos."""
    t = tabla.table
    filas = {t.cell(f, 0).text: f for f in range(3, len(t.rows))}
    tipo = t.cell(filas["1."], 0).fill.fore_color.rgb
    capitulo = t.cell(filas["1.1"], 0).fill.fore_color.rgb
    assert tipo == GRIS_TIPO_DE_COSTE
    assert capitulo == GRIS_SECCION
    assert tipo != capitulo


def test_las_filas_de_datos_no_llevan_fondo(tabla) -> None:
    """Solo las secciones van rellenas: si los datos también lo fueran, no se
    distinguiría el subtotal de la actuación."""
    from pptx.enum.dml import MSO_FILL

    t = tabla.table
    fila_dato = next(f for f in range(3, len(t.rows)) if t.cell(f, 0).text == "1.1.1")
    assert t.cell(fila_dato, 0).fill.type == MSO_FILL.BACKGROUND


def test_la_fila_de_total_va_con_el_tono_fuerte(tabla) -> None:
    t = tabla.table
    ultima = len(t.rows) - 1
    assert t.cell(ultima, 1).text == "TOTAL"
    assert t.cell(ultima, 0).fill.fore_color.rgb == GRIS_TIPO_DE_COSTE


# ─────────────────────────────────────────────────────────────────────────────
#  Lo que se escribe
# ─────────────────────────────────────────────────────────────────────────────


def test_una_actuacion_llega_entera_a_su_fila(tabla) -> None:
    t = tabla.table
    fila = next(f for f in range(3, len(t.rows)) if t.cell(f, 0).text == "1.1.1")
    valores = _texto(t, fila)
    assert valores[1] == "Cubierta"  # objeto
    assert valores[3] == "Cubierta deck al final de su vida útil"
    assert valores[4] == "04 Extremo"
    assert valores[6] == "Vida útil"  # concepto
    assert valores[7] == "NO"  # recuperable
    assert valores[-1] == "412.500,00 €"  # total de la fila
