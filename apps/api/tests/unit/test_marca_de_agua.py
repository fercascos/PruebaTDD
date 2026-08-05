"""Retirada de la marca de agua `[REQ]` P-43."""

from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from tdd.reporting.watermark import hay_marcas_de_agua, retirar_marcas_de_agua


def _presentacion_con(texto: str, *, rotacion: float = 315.0, ancho_in: float = 6.0):
    """Una presentación con una marca de agua en el **patrón**, como la real.

    `MasterShapes` no admite `add_textbox`, así que la forma se crea en una
    diapositiva y se traslada su XML al patrón. Es exactamente donde vive la
    marca en la plantilla del cliente.
    """
    prs = Presentation()
    auxiliar = prs.slides.add_slide(prs.slide_layouts[6])
    forma = auxiliar.shapes.add_textbox(Inches(2), Inches(2), Inches(ancho_in), Inches(2))
    forma.rotation = rotacion
    run = forma.text_frame.paragraphs[0].add_run()
    run.text = texto
    run.font.size = Pt(96)

    elemento = forma._element
    elemento.getparent().remove(elemento)
    prs.slide_masters[0].shapes._spTree.append(elemento)
    return prs


def test_se_detecta_la_marca_del_patron() -> None:
    """Vive en el patrón, no en cada diapositiva: por eso el análisis por
    diapositiva no la vio."""
    assert hay_marcas_de_agua(_presentacion_con("DRAFT")) is True


def test_se_retira_del_patron_y_desaparece_de_todo() -> None:
    prs = _presentacion_con("DRAFT")
    retiradas = retirar_marcas_de_agua(prs)

    assert len(retiradas) == 1
    assert retiradas[0].donde == "patrón"
    assert retiradas[0].texto == "DRAFT"
    assert retiradas[0].rotacion == 315.0
    assert hay_marcas_de_agua(prs) is False


@pytest.mark.parametrize("texto", ["DRAFT", "draft", " Borrador ", "CONFIDENCIAL"])
def test_se_reconocen_las_variantes_habituales(texto: str) -> None:
    prs = _presentacion_con(texto)
    assert len(retirar_marcas_de_agua(prs)) == 1


def test_no_se_borra_un_texto_normal_que_mencione_el_borrador() -> None:
    """La comprobación que evita el destrozo: una diapositiva que hable *sobre*
    el borrador en su cuerpo no puede perder ese texto."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    caja = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    caja.text_frame.text = "Este documento es un borrador y no debe distribuirse"

    assert retirar_marcas_de_agua(prs) == []
    assert "borrador" in caja.text_frame.text


def test_un_titulo_horizontal_que_diga_borrador_no_es_marca_de_agua() -> None:
    """Sin rotación y estrecho: es un título, no una marca de agua."""
    prs = _presentacion_con("BORRADOR", rotacion=0.0, ancho_in=2.0)
    assert retirar_marcas_de_agua(prs) == []


def test_retirarla_dos_veces_no_falla() -> None:
    prs = _presentacion_con("DRAFT")
    assert len(retirar_marcas_de_agua(prs)) == 1
    assert retirar_marcas_de_agua(prs) == []


def test_se_devuelve_lo_retirado_para_poder_auditarlo() -> None:
    """Que el informe emitido haya perdido una forma del patrón no debe ser
    invisible en la auditoría."""
    retiradas = retirar_marcas_de_agua(_presentacion_con("DRAFT"))
    assert retiradas[0].nombre
    assert retiradas[0].donde in ("patrón", "diseño", "diapositiva")
