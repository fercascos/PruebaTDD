"""Retirada de la marca de agua `[REQ]` P-43.

> **P-43 · DECIDIDO por el cliente:** *«la marca de agua de DRAFT no deberá
> aparecer en las versiones futuras, retírala.»*

La marca vive en el **patrón de diapositivas**, no en cada diapositiva: una sola
autoforma rotada 315° con el texto `DRAFT`. Retirarla del patrón la quita de las
67 diapositivas a la vez, que es lo que la hace barata de eliminar y también lo
que la hacía fácil de pasar por alto en un análisis por diapositiva.

`[REQ]` **Esto opera sobre la presentación generada, nunca sobre la plantilla
original.** El original es inmutable: es una de las cuatro garantías del sistema
y hay pruebas que lo verifican por hash.
"""

from __future__ import annotations

from dataclasses import dataclass

from pptx.presentation import Presentation as Presentacion

#: Textos que identifican una marca de agua. Se comparan sin distinguir
#: mayúsculas ni acentos sobrantes, y **solo** contra formas cuyo texto entero
#: sea uno de estos: una diapositiva que hable *sobre* el borrador en su cuerpo
#: no debe perder ese texto.
TEXTOS_DE_MARCA_DE_AGUA = frozenset(
    {"DRAFT", "BORRADOR", "CONFIDENTIAL", "CONFIDENCIAL", "SAMPLE", "MUESTRA"}
)


@dataclass(frozen=True, slots=True)
class MarcaRetirada:
    donde: str  # "patrón" | "diseño" | "diapositiva"
    nombre: str
    texto: str
    rotacion: float


def _es_marca_de_agua(forma) -> bool:  # type: ignore[no-untyped-def]
    if not forma.has_text_frame:
        return False
    texto = forma.text_frame.text.strip().upper()
    if texto not in TEXTOS_DE_MARCA_DE_AGUA:
        return False
    # Una marca de agua está rotada o es muy grande. La comprobación evita
    # borrar un título que legítimamente ponga «BORRADOR» en horizontal.
    rotada = abs(forma.rotation % 360) > 1
    ancha = (forma.width or 0) > 2_700_000  # ~3 in
    return rotada or ancha


def _retirar_de(coleccion, donde: str) -> list[MarcaRetirada]:  # type: ignore[no-untyped-def]
    retiradas: list[MarcaRetirada] = []
    for forma in list(coleccion):
        if _es_marca_de_agua(forma):
            retiradas.append(
                MarcaRetirada(
                    donde=donde,
                    nombre=forma.name or "",
                    texto=forma.text_frame.text.strip(),
                    rotacion=float(forma.rotation or 0),
                )
            )
            forma._element.getparent().remove(forma._element)
    return retiradas


def retirar_marcas_de_agua(prs: Presentacion) -> list[MarcaRetirada]:
    """Quita las marcas de agua de patrones, diseños y diapositivas.

    Devuelve lo retirado, para poder registrarlo: que el informe emitido haya
    perdido una forma del patrón no debe ser invisible en la auditoría.
    """
    retiradas: list[MarcaRetirada] = []
    for patron in prs.slide_masters:
        retiradas += _retirar_de(patron.shapes, "patrón")
        for diseno in patron.slide_layouts:
            retiradas += _retirar_de(diseno.shapes, "diseño")
    for slide in prs.slides:
        retiradas += _retirar_de(slide.shapes, "diapositiva")
    return retiradas


def hay_marcas_de_agua(prs: Presentacion) -> bool:
    """Comprobación de solo lectura, para el análisis de la plantilla."""
    for patron in prs.slide_masters:
        if any(_es_marca_de_agua(f) for f in patron.shapes):
            return True
        for diseno in patron.slide_layouts:
            if any(_es_marca_de_agua(f) for f in diseno.shapes):
                return True
    return any(_es_marca_de_agua(f) for s in prs.slides for f in s.shapes)
