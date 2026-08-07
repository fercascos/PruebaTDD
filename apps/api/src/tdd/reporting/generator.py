"""Generación del informe PPTX `[REQ]` §17.5.

Compone las piezas que ya existían por separado —clonado, sustitución de
marcadores, tabla nativa, inserción de fotos, retirada de la marca de agua— y
las aplica **sobre el snapshot**, nunca sobre la base de datos.

`[REQ]` **La plantilla original no se toca jamás.** Se abre, se trabaja sobre
el objeto en memoria y se guarda en un fichero nuevo. Hay una prueba que
comprueba el hash del fichero de plantilla antes y después.
"""

from __future__ import annotations

import io
from dataclasses import dataclass, field
from decimal import Decimal
from typing import Any

from pptx import Presentation
from pptx.util import Emu, Inches

from tdd.reporting import capex_layout as cl
from tdd.reporting.clone import sustituir_marcadores
from tdd.reporting.pptx_table import insertar_tabla
from tdd.reporting.watermark import retirar_marcas_de_agua

#: Filas de datos por diapositiva de tabla. Medido sobre la plantilla real.
FILAS_POR_DIAPOSITIVA = 18


@dataclass(frozen=True, slots=True)
class FotoParaInsertar:
    """Una foto lista para colocar: sus bytes y su pie."""

    photo_id: str
    datos: bytes
    caption: str = ""


@dataclass
class ResultadoDeGeneracion:
    pptx: bytes
    xlsx: bytes
    diapositivas: int
    diapositivas_de_tabla: int
    marcadores_sin_resolver: list[str] = field(default_factory=list)
    marcas_de_agua_retiradas: list[str] = field(default_factory=list)
    fotos_insertadas: int = 0
    totales: dict[str, Decimal] = field(default_factory=dict)


def valores_de_marcadores(snapshot: dict[str, Any]) -> dict[str, str]:
    """Traduce el snapshot al catálogo cerrado de marcadores de §17.2.

    Se construye **solo desde el snapshot**: si un dato no está congelado, no
    puede salir en el informe, y eso es exactamente lo que se busca.
    """
    proyecto = snapshot["project"]
    activos = snapshot.get("assets", [])
    totales = _totales(snapshot)

    valores = {
        "project.code": str(proyecto.get("internal_code") or ""),
        "project.name": str(proyecto.get("name") or ""),
        "project.client": str(proyecto.get("client_name") or ""),
        "project.currency": str(proyecto.get("currency") or "EUR"),
        "project.asset_count": str(len(activos)),
        "report.generated_at": str(snapshot.get("generated_at", "")),
        "capex.total": cl.formatear_importe(sum(totales.values(), Decimal("0"))),
    }
    for codigo, importe in totales.items():
        valores[f"capex.{codigo.lower()}"] = cl.formatear_importe(importe)
    if activos:
        primero = activos[0]
        valores.update(
            {
                "asset.name": str(primero.get("name") or ""),
                "asset.city": str(primero.get("city") or ""),
                "asset.year_built": str(primero.get("year_built") or ""),
                "asset.total_built_sqm": str(primero.get("total_built_sqm") or ""),
            }
        )
    return valores


def _totales(snapshot: dict[str, Any]) -> dict[str, Decimal]:
    totales: dict[str, Decimal] = {}
    for linea in snapshot.get("capex_items", []):
        codigo = linea["time_horizon_code"]
        totales[codigo] = totales.get(codigo, Decimal("0")) + Decimal(str(linea["amount"]))
    return totales


#: `finding.tenant_recoverable` → lo que escribe la plantilla en esa columna.
#: `NA` se guarda sin punto y se enseña con él, que es como está en la hoja.
RECUPERABLE = {"SI": "SI", "NO": "NO", "NA": "N.A."}


def lineas_de_capex(snapshot: dict[str, Any]) -> list[cl.LineaCapex]:
    """Convierte el snapshot en las líneas que consume `CapexTableLayout`.

    `[REQ]` P-31 · La misma estructura alimenta la tabla nativa del PPTX y la
    hoja del Excel exportado. Construirla dos veces es exactamente lo que hace
    que en seis meses tengan columnas distintas.
    """
    por_hallazgo = {h["id"]: h for h in snapshot.get("findings", [])}
    salida: list[cl.LineaCapex] = []
    for linea in snapshot.get("capex_items", []):
        hallazgo = por_hallazgo.get(linea["finding_id"])
        if hallazgo is None:
            continue
        salida.append(
            cl.LineaCapex(
                numero="",
                zona=str(hallazgo.get("zone_name") or ""),
                concepto=str(hallazgo.get("concept_name") or ""),
                descripcion=str(hallazgo.get("title") or ""),
                riesgo=str(hallazgo.get("risk_name") or ""),
                comentarios=str(hallazgo.get("comments") or hallazgo.get("description") or ""),
                horizonte=str(linea["time_horizon_code"]),
                importe=Decimal(str(linea["amount"])),
                # Los tres niveles del árbol vienen ya resueltos del snapshot.
                # `capex_item_name` solo trae valor si el hallazgo apunta al
                # tercer nivel: un hallazgo codificado en el capítulo se queda
                # sin objeto, y sale con la celda vacía en vez de repetir el
                # nombre del capítulo, que sería mentir por partida doble.
                tipo_de_coste=str(hallazgo.get("capex_type_name") or ""),
                capitulo=str(hallazgo.get("capex_chapter_name") or ""),
                objeto=str(hallazgo.get("capex_item_name") or ""),
                recuperable=RECUPERABLE.get(str(hallazgo.get("tenant_recoverable") or ""), ""),
                # [REQ] P-44 · Las líneas de la misma actuación se agrupan en
                # una sola fila con varias columnas de plazo rellenas.
                finding_id=str(hallazgo["id"]),
            )
        )
    return salida


def insertar_foto(
    slide: Any, datos: bytes, *, izquierda: float, arriba: float, ancho_max: float, alto_max: float
) -> Any:
    """Coloca una imagen **conservando su proporción** dentro del hueco.

    Estirar una foto para que llene el marco es el error más visible que puede
    tener un informe técnico: un edificio deformado desacredita el documento
    entero antes de que nadie lea una línea.
    """
    from PIL import Image

    with Image.open(io.BytesIO(datos)) as img:
        ancho_px, alto_px = img.size
    proporcion = ancho_px / alto_px if alto_px else 1.0

    ancho = ancho_max
    alto = ancho / proporcion
    if alto > alto_max:
        alto = alto_max
        ancho = alto * proporcion

    # Centrado en el hueco: lo que sobra se reparte, no se deja todo a un lado.
    desplazamiento_x = (ancho_max - ancho) / 2
    desplazamiento_y = (alto_max - alto) / 2
    return slide.shapes.add_picture(
        io.BytesIO(datos),
        Inches(izquierda + desplazamiento_x),
        Inches(arriba + desplazamiento_y),
        Inches(ancho),
        Inches(alto),
    )


def _hueco_libre(
    slide: Any, ancho_diapositiva: float, alto_diapositiva: float
) -> tuple[float, float, float, float]:
    """Un hueco razonable para las fotos: la mitad inferior de la diapositiva."""
    margen = 0.6
    return (
        margen,
        alto_diapositiva / 2,
        ancho_diapositiva - 2 * margen,
        alto_diapositiva / 2 - margen,
    )


def generar(
    plantilla: bytes,
    snapshot: dict[str, Any],
    *,
    fotos: list[FotoParaInsertar] | None = None,
    locale: str = "es-ES",
    retirar_marca_de_agua: bool = True,
) -> ResultadoDeGeneracion:
    """Produce el PPTX y el XLSX. **Sobre bytes: el original no se toca.**"""
    from tdd.exports.capex_xlsx import generar_xlsx

    # `Presentation` sobre BytesIO: la plantilla del disco no se abre ni se
    # toca, se trabaja sobre una copia en memoria.
    prs = Presentation(io.BytesIO(plantilla))

    # 1 · Marcadores, en todas las diapositivas que los tengan.
    valores = valores_de_marcadores(snapshot)
    sin_resolver: list[str] = []
    for slide in prs.slides:
        sin_resolver += sustituir_marcadores(slide, valores)

    # 2 · Tabla nativa de CAPEX, partida si hace falta.
    layout = cl.construir(lineas_de_capex(snapshot), capitulo="CAPEX", locale=locale)
    trozos = cl.particionar(layout, filas_por_diapositiva=FILAS_POR_DIAPOSITIVA)
    for trozo in trozos:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        insertar_tabla(slide, trozo)

    # 3 · Fotografías, en el orden que fijó el consultor.
    insertadas = 0
    if fotos:
        ancho_in = Emu(prs.slide_width).inches
        alto_in = Emu(prs.slide_height).inches
        for foto in fotos:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            izquierda, arriba, ancho_max, alto_max = _hueco_libre(slide, ancho_in, alto_in)
            try:
                insertar_foto(
                    slide,
                    foto.datos,
                    izquierda=izquierda,
                    arriba=arriba,
                    ancho_max=ancho_max,
                    alto_max=alto_max,
                )
            except Exception:  # noqa: BLE001 — una foto ilegible no tumba el informe
                continue
            if foto.caption:
                caja = slide.shapes.add_textbox(
                    Inches(izquierda), Inches(arriba + alto_max), Inches(ancho_max), Inches(0.35)
                )
                caja.text_frame.text = foto.caption
            insertadas += 1

    # 4 · [REQ] P-43 · La marca de agua de borrador no aparece en lo generado.
    marcas = retirar_marcas_de_agua(prs) if retirar_marca_de_agua else []

    salida = io.BytesIO()
    prs.save(salida)
    return ResultadoDeGeneracion(
        pptx=salida.getvalue(),
        xlsx=generar_xlsx(layout),
        diapositivas=len(prs.slides),
        diapositivas_de_tabla=len(trozos),
        marcadores_sin_resolver=sorted(set(sin_resolver)),
        marcas_de_agua_retiradas=[m.texto for m in marcas],
        fotos_insertadas=insertadas,
        totales=layout.totales,
    )


def analizar(plantilla: bytes) -> dict[str, Any]:
    """Lee la plantilla y devuelve lo que hace falta para mapearla `[REQ]` §17.3.

    No modifica nada: es una lectura. Lo que devuelve alimenta la pantalla de
    mapeo, y sin ella el usuario tendría que adivinar qué marcadores existen.
    """
    import re

    from tdd.reporting.clone import texto_completo
    from tdd.reporting.watermark import hay_marcas_de_agua

    prs = Presentation(io.BytesIO(plantilla))
    patron = re.compile(r"\{\{\s*([a-zA-Z0-9_.]+)\s*\}\}")

    marcadores: set[str] = set()
    por_diapositiva: list[dict[str, Any]] = []
    for indice, slide in enumerate(prs.slides, start=1):
        texto = texto_completo(slide)
        encontrados = sorted(set(patron.findall(texto)))
        marcadores.update(encontrados)
        por_diapositiva.append(
            {
                "index": indice,
                "placeholders": encontrados,
                "shapes": len(slide.shapes),
                "has_table": any(f.has_table for f in slide.shapes),
                "has_picture": any(f.shape_type == 13 for f in slide.shapes),  # noqa: PLR2004
            }
        )

    fuentes = sorted(
        {
            run.font.name
            for slide in prs.slides
            for forma in slide.shapes
            if forma.has_text_frame
            for parrafo in forma.text_frame.paragraphs
            for run in parrafo.runs
            if run.font.name
        }
    )
    return {
        "slide_count": len(prs.slides),
        "slide_width_in": round(Emu(prs.slide_width).inches, 2),
        "slide_height_in": round(Emu(prs.slide_height).inches, 2),
        "placeholders": sorted(marcadores),
        "slides": por_diapositiva,
        "fonts": fuentes,
        # [REQ] P-43 · Se avisa de que la plantilla la trae, porque la
        # generación la retirará y conviene que nadie se sorprenda.
        "has_watermark": hay_marcas_de_agua(prs),
    }
