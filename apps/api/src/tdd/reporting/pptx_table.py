"""Tabla nativa de CAPEX en PowerPoint `[REQ]` P-31.

Sustituye a la imagen EMF pegada desde Excel. Los colores y anchos salen del
render de la plantilla real del cliente, no de una propuesta.
"""

from __future__ import annotations

from typing import Any

from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.slide import Slide
from pptx.util import Emu, Inches, Pt

from tdd.reporting.capex_layout import TITULO_GRUPO, Alineacion, CapexTableLayout

#: Muestreados del render de la plantilla del cliente (docs/20 §20.3).
VERDE_TITULO = RGBColor(0xA9, 0xC7, 0x8C)
GRIS_CABECERA = RGBColor(0xB0, 0xB0, 0xB0)
ORO_GRUPO = RGBColor(0x9A, 0x8C, 0x4E)
GRIS_SECCION = RGBColor(0xA6, 0xA6, 0xA6)
#: Las secciones de nivel 1 —el tipo de coste— y la fila de TOTAL, más oscuras
#: que las de capítulo para que se lean como lo que son: el grupo de arriba.
GRIS_TIPO_DE_COSTE = RGBColor(0x59, 0x59, 0x59)
BLANCO = RGBColor(0xFF, 0xFF, 0xFF)
NEGRO = RGBColor(0x00, 0x00, 0x00)

#: Un color por plazo, como en el original.
COLOR_PLAZO = {
    "corto": RGBColor(0xF8, 0xCB, 0xCB),
    "medio": RGBColor(0xFB, 0xE5, 0xA6),
    "largo": RGBColor(0xC8, 0xE6, 0xC9),
    "mejoras": RGBColor(0xBD, 0xD7, 0xEE),
    "otro": RGBColor(0xE0, 0xE0, 0xE0),
}

#: P-38 · toda la tipografía unificada en Gotham. El original mezclaba Century
#: Gothic dentro de las imágenes de tabla, porque venían de un Excel ajeno.
FUENTE_CUERPO = "Gotham Light"
FUENTE_CABECERA = "Gotham Medium"

_ALIGN = {
    Alineacion.IZQUIERDA: PP_ALIGN.LEFT,
    Alineacion.CENTRO: PP_ALIGN.CENTER,
    Alineacion.DERECHA: PP_ALIGN.RIGHT,
}


def _escribir(
    celda: Any, texto: str, *, pt: float, negrita: bool, color: Any, alineacion: Any
) -> None:
    celda.text_frame.word_wrap = True
    p = celda.text_frame.paragraphs[0]
    p.alignment = alineacion
    run = p.add_run()
    run.text = texto
    run.font.size = Pt(pt)
    run.font.bold = negrita
    run.font.color.rgb = color
    run.font.name = FUENTE_CABECERA if negrita else FUENTE_CUERPO
    celda.margin_left = celda.margin_right = Emu(27432)  # 0,03 in
    celda.margin_top = celda.margin_bottom = Emu(9144)


def insertar_tabla(
    slide: Slide,
    layout: CapexTableLayout,
    *,
    left_in: float = 0.47,
    top_in: float = 1.55,
    alto_fila_in: float = 0.17,
    cuerpo_pt: float = 5.0,
) -> None:
    """Dibuja la tabla en la diapositiva, con su cabecera de dos niveles."""
    n_col = len(layout.columnas)
    # 2 filas de cabecera de la tabla + 1 de título de bloque
    n_fil = len(layout.filas) + 3

    forma = slide.shapes.add_table(
        n_fil,
        n_col,
        Inches(left_in),
        Inches(top_in),
        Inches(layout.ancho_total_in),
        Inches(alto_fila_in * n_fil),
    )
    tabla = forma.table
    for i, c in enumerate(layout.columnas):
        tabla.columns[i].width = Inches(c.ancho_in)
    for f in range(n_fil):
        tabla.rows[f].height = Inches(alto_fila_in)

    # ── Fila 0 · título del bloque, combinado de lado a lado ────────────────
    tabla.cell(0, 0).merge(tabla.cell(0, n_col - 1))
    c0 = tabla.cell(0, 0)
    c0.fill.solid()
    c0.fill.fore_color.rgb = VERDE_TITULO
    _escribir(
        c0, layout.titulo, pt=cuerpo_pt + 2, negrita=True, color=BLANCO, alineacion=PP_ALIGN.CENTER
    )

    # ── Filas 1-2 · cabecera de dos niveles ─────────────────────────────────
    primera = next(i for i, c in enumerate(layout.columnas) if c.grupo == "capex")
    ultima = max(i for i, c in enumerate(layout.columnas) if c.grupo == "capex")

    for i, col in enumerate(layout.columnas):
        if col.grupo is None:
            # Sin grupo: la cabecera ocupa las dos filas.
            tabla.cell(1, i).merge(tabla.cell(2, i))
            celda = tabla.cell(1, i)
            celda.fill.solid()
            celda.fill.fore_color.rgb = ORO_GRUPO if col.key == "riesgo" else GRIS_CABECERA
            _escribir(
                celda,
                layout.titulo_columna(col),
                pt=cuerpo_pt,
                negrita=True,
                color=BLANCO if col.key == "riesgo" else NEGRO,
                alineacion=PP_ALIGN.CENTER,
            )
        else:
            celda = tabla.cell(2, i)
            celda.fill.solid()
            celda.fill.fore_color.rgb = COLOR_PLAZO.get(col.key, GRIS_CABECERA)
            _escribir(
                celda,
                layout.titulo_columna(col),
                pt=cuerpo_pt,
                negrita=False,
                color=NEGRO,
                alineacion=PP_ALIGN.CENTER,
            )

    tabla.cell(1, primera).merge(tabla.cell(1, ultima))
    cg = tabla.cell(1, primera)
    cg.fill.solid()
    cg.fill.fore_color.rgb = ORO_GRUPO
    _escribir(
        cg,
        TITULO_GRUPO["capex"][0 if layout.locale.startswith("es") else 1],
        pt=cuerpo_pt,
        negrita=True,
        color=BLANCO,
        alineacion=PP_ALIGN.CENTER,
    )

    # ── Cuerpo ──────────────────────────────────────────────────────────────
    #
    # Las secciones vienen en dos niveles: el tipo de coste —«HARD COSTS»— y
    # dentro el capítulo. Se distinguen por tono, como en la hoja del cliente:
    # si los dos fuesen del mismo gris, una tabla con cinco capítulos parecería
    # tener diez secciones sueltas en vez de dos grupos.
    for f, fila in enumerate(layout.filas, start=3):
        seccion = fila.tipo in ("seccion", "total")
        fondo = GRIS_SECCION if fila.nivel == 2 else GRIS_TIPO_DE_COSTE
        for i, col in enumerate(layout.columnas):
            celda = tabla.cell(f, i)
            if seccion:
                celda.fill.solid()
                celda.fill.fore_color.rgb = fondo
            else:
                celda.fill.background()
            _escribir(
                celda,
                fila.celdas.get(col.key, ""),
                pt=cuerpo_pt,
                negrita=seccion,
                color=BLANCO if seccion else NEGRO,
                alineacion=_ALIGN[col.alineacion],
            )
