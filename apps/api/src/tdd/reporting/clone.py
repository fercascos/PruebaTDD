"""Clonado de diapositivas y sustitución de marcadores.

`[LIM]` **`python-pptx` no ofrece duplicado oficial de diapositivas.** Se clona
copiando el XML del árbol de formas y volviendo a enlazar las relaciones. Es la
técnica conocida, y funciona; lo que no hace es garantizarse contra cualquier
plantilla arbitraria. El corpus de pruebas la ejerce contra las cuatro reales.
"""

from __future__ import annotations

import copy
import re
from collections.abc import Iterator
from typing import Any

from pptx.presentation import Presentation as Presentacion
from pptx.slide import Slide

MARCADOR = re.compile(r"\{\{([a-zA-Z0-9_.|:# ]+)\}\}")


def clonar_diapositiva(prs: Presentacion, origen: Slide) -> Slide:
    """Duplica una diapositiva al final de la presentación."""
    destino = prs.slides.add_slide(origen.slide_layout)

    # La diapositiva nueva viene con los marcadores de posición del diseño;
    # se retiran porque el contenido real llega copiado del origen.
    for forma in list(destino.shapes):
        forma._element.getparent().remove(forma._element)

    for forma in origen.shapes:
        destino.shapes._spTree.append(copy.deepcopy(forma._element))

    # Reenlazar las partes referenciadas —imágenes, sobre todo—. Sin esto, las
    # fotos del origen apuntan a relaciones que la diapositiva nueva no tiene y
    # PowerPoint da el fichero por dañado.
    #
    # Los identificadores NO se pueden reutilizar: la diapositiva nueva ya trae
    # los suyos (empezando por el enlace a su diseño) y colisionarían. Se
    # construye un mapa antiguo → nuevo y se reescriben las referencias del XML
    # ya copiado.
    mapa: dict[str, str] = {}
    for rid, rel in origen.part.rels.items():
        if rel.reltype.endswith("slideLayout"):
            continue  # la diapositiva nueva ya tiene el suyo
        if rel.is_external:
            mapa[rid] = destino.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        else:
            mapa[rid] = destino.part.relate_to(rel.target_part, rel.reltype)

    if mapa:
        _reescribir_referencias(destino, mapa)
    resultado: Slide = destino
    return resultado


#: Atributos que llevan un identificador de relación en el XML de una forma.
_ATRIBUTOS_RID = (
    "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed",
    "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}link",
    "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id",
)


def _reescribir_referencias(slide: Slide, mapa: dict[str, str]) -> None:
    for elemento in slide.shapes._spTree.iter():
        for atributo in _ATRIBUTOS_RID:
            valor = elemento.get(atributo)
            if valor in mapa:
                elemento.set(atributo, mapa[valor])


def _runs_de(forma: Any) -> Iterator[Any]:
    if not forma.has_text_frame:
        return
    yield from forma.text_frame.paragraphs


def sustituir_marcadores(slide: Slide, valores: dict[str, str]) -> list[str]:
    """Sustituye `{{clave}}` conservando el formato. Devuelve los no resueltos.

    PowerPoint parte un marcador entre varios `run` con pasmosa facilidad —basta
    con que alguien corrigiera una letra al escribirlo—. Por eso se opera sobre
    el texto del párrafo completo y se vuelca en el primer `run`, que es el que
    lleva el formato bueno.
    """
    sin_resolver: list[str] = []
    for forma in slide.shapes:
        for parrafo in _runs_de(forma):
            runs = parrafo.runs
            if not runs:
                continue
            texto = "".join(r.text for r in runs)
            if "{{" not in texto:
                continue

            def _rep(m: re.Match[str]) -> str:
                clave = m.group(1).strip()
                if clave in valores:
                    return valores[clave]
                sin_resolver.append(clave)
                # Se deja vacío, nunca el literal: un `{{...}}` impreso en un
                # informe entregado al cliente es el peor resultado posible.
                return ""

            nuevo = MARCADOR.sub(_rep, texto)
            runs[0].text = nuevo
            for r in runs[1:]:
                r.text = ""
    return sin_resolver


def texto_completo(slide: Slide) -> str:
    partes = []
    for forma in slide.shapes:
        if forma.has_text_frame:
            partes.append(forma.text_frame.text)  # type: ignore[attr-defined]
        if getattr(forma, "has_table", False):
            for fila in forma.table.rows:  # type: ignore[attr-defined]
                for celda in fila.cells:
                    partes.append(celda.text)
    return "\n".join(partes)
